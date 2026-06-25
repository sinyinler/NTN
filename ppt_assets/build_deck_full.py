# -*- coding: utf-8 -*-
"""构建 NTN 复现【全流程】汇报 PPT：每页讲 做了什么 / 为什么这么试 / 结果。"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from PIL import Image

A = r'D:\Desktop\NTN\ppt_assets'
R = r'D:\Desktop\results'
OUT = r'D:\Desktop\NTN\NTN_复现全流程汇报.pptx'

NAVY=RGBColor(0x21,0x29,0x5C); BLUE=RGBColor(0x06,0x5A,0x82); TEAL=RGBColor(0x1C,0x72,0x93)
MINT=RGBColor(0x00,0xA8,0x96); INK=RGBColor(0x2A,0x2E,0x35); MUTE=RGBColor(0x6B,0x76,0x82)
WHITE=RGBColor(0xFF,0xFF,0xFF); SOFT=RGBColor(0xF2,0xF5,0xF7); LINE=RGBColor(0xD9,0xE0,0xE5)
RED=RGBColor(0xC0,0x56,0x3A); ICE=RGBColor(0xCA,0xDC,0xFC)
TITLEF="Microsoft YaHei"; BODYF="Microsoft YaHei"

prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
BLANK=prs.slide_layouts[6]; SW,SH=13.333,7.5

def slide(bg=WHITE):
    s=prs.slides.add_slide(BLANK)
    r=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,prs.slide_width,prs.slide_height)
    r.fill.solid(); r.fill.fore_color.rgb=bg; r.line.fill.background(); r.shadow.inherit=False
    s.shapes._spTree.remove(r._element); s.shapes._spTree.insert(2,r._element)
    return s

def box(s,l,t,w,h,fill=None,line=None,round_=False,lw=1.0):
    shp=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if round_ else MSO_SHAPE.RECTANGLE,Inches(l),Inches(t),Inches(w),Inches(h))
    if fill is None: shp.fill.background()
    else: shp.fill.solid(); shp.fill.fore_color.rgb=fill
    if line is None: shp.line.fill.background()
    else: shp.line.color.rgb=line; shp.line.width=Pt(lw)
    shp.shadow.inherit=False; return shp

def text(s,l,t,w,h,runs,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,space=1.06):
    tb=s.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h)); tf=tb.text_frame
    tf.word_wrap=True; tf.vertical_anchor=anchor
    tf.margin_left=0; tf.margin_right=0; tf.margin_top=0; tf.margin_bottom=0
    if isinstance(runs[0],tuple): runs=[runs]
    for i,para in enumerate(runs):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.alignment=align; p.line_spacing=space; p.space_after=Pt(4)
        for (txt,sz,col,bold,*rest) in para:
            r=p.add_run(); r.text=txt; f=r.font
            f.size=Pt(sz); f.color.rgb=col; f.bold=bold; f.name=(rest[0] if rest else BODYF)
    return tb

def motif_title(s,title,num,dark=False):
    box(s,0.6,0.5,0.24,0.24,fill=MINT)
    text(s,0.98,0.4,11.0,0.7,[[(title,27,WHITE if dark else NAVY,True,TITLEF)]])
    text(s,12.2,0.46,0.7,0.5,[[("%02d"%num,13,MUTE,True,TITLEF)]],align=PP_ALIGN.RIGHT)

def pic(s,path,l,t,w,h):
    iw,ih=Image.open(path).size; a=iw/ih; ba=w/h
    nw,nh=(w,w/a) if a>ba else (h*a,h)
    s.shapes.add_picture(path,Inches(l+(w-nw)/2),Inches(t+(h-nh)/2),Inches(nw),Inches(nh))

def bullets(s,l,t,w,h,items,sz=14.5,gap=7,col=INK,marker=MINT):
    tb=s.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h)); tf=tb.text_frame
    tf.word_wrap=True; tf.margin_left=0; tf.margin_top=0; tf.margin_right=0
    for i,it in enumerate(items):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.line_spacing=1.12; p.space_after=Pt(gap)
        r=p.add_run(); r.text="▪  "; r.font.size=Pt(sz); r.font.color.rgb=marker; r.font.bold=True; r.font.name=BODYF
        if isinstance(it,str): it=[(it,sz,col,False)]
        for (txt,s2,c2,b2) in it:
            r=p.add_run(); r.text=txt; r.font.size=Pt(s2); r.font.color.rgb=c2; r.font.bold=b2; r.font.name=BODYF
    return tb

def why(s,l,t,w,txt):  # 统一的"为什么"小条
    box(s,l,t,w,0.62,fill=RGBColor(0xEC,0xF6,0xF4),round_=True)
    text(s,l+0.2,t+0.09,w-0.4,0.5,[[("为什么这么试： ",13,MINT,True),(txt,13,INK,False)]],anchor=MSO_ANCHOR.MIDDLE,space=1.05)

def cap(s,l,t,w,txt):
    text(s,l,t,w,0.35,[[(txt,10.5,MUTE,False)]],align=PP_ALIGN.CENTER)

# ============ 1 Title ============
s=slide(NAVY); box(s,0,0,0.18,SH,fill=MINT)
text(s,1.1,2.0,11.2,2.2,[[("把噪声“翻译”成高斯：NTN 用于 BFI 去噪的",34,WHITE,True,TITLEF)],
                         [("复现全流程与泛化性研究",34,MINT,True,TITLEF)]],space=1.12)
text(s,1.12,4.3,11,1.0,[[("复现并改造《Learning to Translate Noise for Robust Image Denoising》(arXiv:2412.04727)",15,ICE,False)],
                        [("到无真值 GT 的 BFI / Noise2Noise 场景 —— 含诊断、修复、机制验证、多被试泛化与边界分析",15,ICE,False)]])
text(s,1.12,6.5,11,0.5,[[("宋亚栋  ·  组会汇报（全流程）  ·  2026-06",12.5,MUTE,False)]])

# ============ 2 背景 ============
s=slide(); motif_title(s,"背景：BFI 去噪的“泛化”难题",2)
bullets(s,0.7,1.55,7.0,4.6,[
    [("BFI（血流成像）原始图噪声大，",15,INK,False),("去噪是后续分析的必要前处理。",15,INK,False)],
    [("深度去噪器在“训练见过的噪声”上很好，",15,INK,False),("换采集设置/部位/噪声分布就明显变差。",15,RED,True)],
    [("两类 OOD 要区分：",15,NAVY,True)],
    [("  · 噪声 OOD（不同采集/噪声强度，如 3×3 vs 5×5 窗口）",14,INK,False)],
    [("  · 内容 OOD（没见过的解剖结构，如鼠脑训练→人脚测试）",14,INK,False)],
    [("目标：在有限、无真值 GT 的条件下，做到对 OOD 也稳。",15,INK,False)],
],sz=15,gap=11)
box(s,8.1,2.0,4.5,1.3,fill=SOFT,line=LINE,round_=True)
text(s,8.35,2.2,4.0,1.0,[[("训练见过的噪声",13,MUTE,False)],[("去噪好  ✔",18,TEAL,True)]],space=1.1)
box(s,8.1,3.9,4.5,1.3,fill=SOFT,line=LINE,round_=True)
text(s,8.35,4.1,4.0,1.0,[[("没见过的噪声/结构 (OOD)",13,MUTE,False)],[("残留 / 过平滑 / 脑补  ✘",17,RED,True)]],space=1.1)
ar=s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW,Inches(10.1),Inches(3.32),Inches(0.45),Inches(0.5)); ar.fill.solid(); ar.fill.fore_color.rgb=MUTE; ar.line.fill.background(); ar.shadow.inherit=False

# ============ 3 定位 ============
s=slide(); motif_title(s,"方法定位：与已有方法的区别",3)
cards=[("合成高斯 / 监督去噪","只认高斯噪声；遇真实噪声崩。",TEAL),
       ("N2N / 真实数据训练","受限训练噪声分布，OOD 易过拟合。",BLUE),
       ("NTN（本文）","不直接去未知噪声，先“翻译”成高斯，再用通用高斯去噪器；T 与去噪器解耦、即插即用。",MINT)]
cw=3.95; x0=0.7
for i,(h,b,c) in enumerate(cards):
    x=x0+i*(cw+0.28); box(s,x,1.9,cw,3.9,fill=SOFT,line=LINE,round_=True); box(s,x,1.9,cw,0.12,fill=c)
    text(s,x+0.3,2.2,cw-0.6,0.9,[[(h,17,c,True)]]); text(s,x+0.3,3.05,cw-0.6,2.6,[[(b,14.5,INK,False)]],space=1.16)
text(s,0.7,6.1,12,0.5,[[("一句话：把“去各种噪声”的难题，转化为“把各种噪声翻译成同一种高斯”。",13.5,MUTE,True)]])

# ============ 4 核心思路 ============
s=slide(); motif_title(s,"核心思路：与其硬去噪，不如先“翻译”噪声",4)
text(s,0.7,1.5,11.9,1.3,[
    [("关键观察（论文 Fig.1）：",14.5,NAVY,True),("给真实噪声图再加点高斯反而更好去——去噪器最擅长高斯。",14.5,INK,False)],
    [("于是先用小网络 ",14.5,INK,False),("T 把未知噪声翻译成空间不相关的高斯噪声",14.5,MINT,True),("，再交给高斯去噪器 D。",14.5,INK,False)]],space=1.12)
pic(s,os.path.join(A,'fig1_teaser.png'),0.7,3.0,11.9,3.5)
cap(s,0.7,6.55,11.9,"论文 Fig.1：OOD 噪声图 → 翻译网络 → 高斯噪声图 → 高斯去噪器 → 干净图")

# ============ 5 方法框架 ============
s=slide(); motif_title(s,"方法框架：T 翻译 + D 去高斯（两阶段两损失）",5)
pic(s,os.path.join(A,'fig3_arch.png'),0.6,1.5,7.3,3.0); cap(s,0.6,4.55,7.3,"论文 Fig.3：T 可训练🔥，D 冻结❄")
bullets(s,8.2,1.6,4.6,5.0,[
    [("阶段一：",14.5,NAVY,True),("训练只会去高斯噪声的 D。",14.5,INK,False)],
    [("阶段二：",14.5,NAVY,True),("冻结 D，训 T 使 D(T(I))≈干净。",14.5,INK,False)],
    [("implicit loss：",14.5,NAVY,True),("去噪结果像干净图。",14.5,INK,False)],
    [("explicit loss：",14.5,NAVY,True),("Wasserstein 把翻译后噪声逼成高斯：空间域正态 + 频域 Rayleigh(更白)。",14.5,INK,False)],
    [("泛化来源：",14.5,MINT,True),("D 只认高斯(不过拟合真实噪声)，T 把任意噪声搬到高斯。",14.5,INK,False)],
],sz=14.5,gap=12)

# ============ 6 诊断起点 ============
s=slide(); motif_title(s,"复现起点：原项目为什么“效果差、无泛化”",6)
why(s,0.7,1.5,12.0,"对照论文逐行审查代码，先定位病根，而不是盲目调参。")
bullets(s,0.7,2.45,12.0,3.8,[
    [("最致命：D' 被做成“盲高斯 + σ∈[0,0.15]”，而非论文那台固定-σ偏科去噪器。",15,RED,True)],
    [("  · σ 下界到 0 → “什么都不动”也是合法解 → T 学成恒等映射、翻译被短路；",14,INK,False)],
    [("  · σ 上界仅 0.15 → 后面实测真实噪声 0.10~0.43，D' 根本没见过强噪 → 推理必崩。",14,INK,False)],
    [("implicit 用 I2、explicit 用糊均值当锚，且 explicit 把“血管结构”当噪声去高斯化 → 磨血管。",14.5,INK,False)],
    [("explicit 从第 0 步全程开（论文应后 50%）；lr 偏高；GIBlock 注入初值为 0 形同虚设。",14.5,INK,False)],
],sz=14.5,gap=9)

# ============ 7 测噪声+重构 ============
s=slide(); motif_title(s,"第一步：实测真实噪声 → 据此重构",7)
why(s,0.7,1.5,12.0,"D' 该对齐的 σ 必须等于“真实噪声在 log1p 域的强度”，否则训练/推理对不上。")
text(s,0.7,2.35,5.6,0.4,[[("实测 5x5（log1p 域 σ，相邻帧差/MAD）：",13.5,NAVY,True)]])
rows=[("叠加层级","σ_mad","~1/N"),("5x5x1","0.43","1.00"),("5x5x2","0.21","0.50"),("5x5x3","0.14","0.33"),("5x5x4","0.10","0.25")]
yy=2.8
for i,(a_,b_,c_) in enumerate(rows):
    bg=SOFT if i==0 else WHITE; box(s,0.7,yy,5.2,0.46,fill=bg,line=LINE)
    for cx,txt in zip([0.9,3.0,4.4],(a_,b_,c_)):
        text(s,cx,yy+0.07,1.6,0.4,[[(txt,12.5,(NAVY if i==0 else INK),i==0)]])
    yy+=0.46
bullets(s,6.4,2.5,6.3,4.2,[
    [("噪声近完美按 1/N 衰减 → 测量可信、跨度达 4 倍。",14,INK,False)],
    [("据此重构（每步一个 commit）：",14.5,NAVY,True)],
    [("  · D' 改“盲区间 σ∈[0.08,0.6]”覆盖真实跨度、下界离 0 堵恒等捷径；",13.5,INK,False)],
    [("  · Ĉ 改用 N2N(I1)（内容对齐，翻译噪声才“纯”）；",13.5,INK,False)],
    [("  · explicit loss 后 50% 才启用；lr→1e-3 cosine；GIBlock 注入生效。",13.5,INK,False)],
],sz=14,gap=8)

# ============ 8 公平对照 ============
s=slide(); motif_title(s,"实验设计：公平对照 + 留出 OOD",8)
why(s,0.7,1.5,12.0,"要证明泛化，N2N 基线必须与 NTN 用完全相同的数据；并留出没训过的层级作 OOD。")
bullets(s,0.7,2.5,12.0,3.4,[
    [("N2N 与 NTN 都只用 level 2/3/4 训练；最噪的 level1 完全留作 OOD 测试。",15,INK,False)],
    [("N2N 既是 NTN 的 Ĉ 生成器，又是对照基线 → 唯一差别只剩“NTN 翻译框架”。",15,INK,False)],
    [("评判标准：不只看 PSNR/SSIM/r，更看血管局部放大（是否过平滑/脑补）。",15,MINT,True)],
    [("（修复隐患：发现数据发现逻辑会漏掉混合结构场景；按叠加层级筛选 --levels。）",13.5,MUTE,False)],
],sz=15,gap=12)

# ============ 9 结果① level1 OOD ============
s=slide(); motif_title(s,"结果①：level1（OOD）定量，NTN 略胜 N2N",9)
pic(s,os.path.join(A,'chart_psnr.png'),0.5,1.7,5.6,4.2)
pic(s,os.path.join(R,'eval_ood','compare','scene0_frame0.png'),6.2,1.7,6.7,3.6)
cap(s,6.2,5.35,6.7,"noisy / N2N / NTN / 伪GT(level4多帧均值)，上全图下放大")
text(s,6.2,5.8,6.7,0.9,[[("NTN 25.64 vs N2N 25.01 dB（+0.63 dB），SSIM 持平；",14,INK,False)],
                        [("39 场景、log1p 域、同数据训练同 GT 评测。",13,MUTE,False)]],space=1.1)

# ============ 10 机制验证 ============
s=slide(); motif_title(s,"结果②：翻译机制成立（硬证据）",10)
why(s,0.7,1.5,12.0,"要确认 NTN 的好处来自“真的把噪声翻译成白高斯”，而非偶然——否则结论站不住。")
pic(s,os.path.join(A,'chart_decorr.png'),0.5,2.35,5.0,3.0)
pic(s,os.path.join(R,'diag','noise_panels','scene0_frame0.png'),5.7,2.35,7.1,3.1)
cap(s,5.7,5.5,7.1,"I / T(I) / |真实噪声| / |翻译后噪声|：后者明显更细更白")
text(s,0.5,5.55,5.0,1.2,[[("真实散斑空间相关(自相关0.82)→翻译后近白(0.13)，谱平坦度×9；",12.5,INK,False)],
                        [("D'(I)直接去=14dB(崩) → D'(T(I))=20dB(+6)，翻译是“救命”一步。",12.5,INK,False)]],space=1.08)

# ============ 11 内容OOD问题 ============
s=slide(); motif_title(s,"新问题：内容 OOD —— 人手出现“脑补”假血管",11)
why(s,0.7,1.5,12.0,"NTN 解决“噪声 OOD”，但解决不了“没见过的解剖结构”——去噪器套用训练血管先验。")
bullets(s,0.7,2.5,12.0,2.2,[
    [("在没见过的人手 BFI 上，模型凭空生成训练里那种密集细血管（hallucination）。",15,RED,True)],
    [("根因：去噪器的图像先验偏向训练内容（鼠脑/皮层），低信噪区用先验“填”出血管。",14.5,INK,False)],
    [("对策：把多种被试加进训练，让模型见过多样的血管形态。",15,MINT,True)],
])
box(s,0.7,5.0,12.0,1.4,fill=SOFT,line=LINE,round_=True)
text(s,1.0,5.2,11.4,1.0,[[("关键区分： ",14,NAVY,True),
    ("NTN(翻译) 管“噪声没见过”；“结构没见过”要靠训练内容多样化（喂强 N2N/D'）。两者是不同维度的 OOD。",14,INK,False)]],space=1.12,anchor=MSO_ANCHOR.MIDDLE)

# ============ 12 多被试 + 脚 老vs新 ============
s=slide(); motif_title(s,"结果③：多被试训练“救活”了没见过的人脚",12)
text(s,0.7,1.42,12.0,0.5,[[("加入 鼠脑+鼠腿+人手 多被试重训（σ→[0.02,0.6]，N2N 同数据重训）；脚完全留作 OOD。",13.5,INK,False)]])
pic(s,os.path.join(R,'metrics_foot5_old','compare_gray_color.png'),0.6,2.0,6.2,2.3)
text(s,0.6,4.32,6.2,0.4,[[("旧模型(只5x5)：NTN 全红、r = −0.15 → 崩溃",12.5,RED,True)]],align=PP_ALIGN.CENTER)
pic(s,os.path.join(R,'images','foot_5x5x5_50','comparison','5x5x5_0_npy_50_input_n2n_ntn.png'),7.0,2.0,5.9,2.3)
text(s,7.0,4.32,5.9,0.4,[[("新模型(多被试)：正常去噪、不崩不脑补",12.5,MINT,True)]],align=PP_ALIGN.CENTER)
box(s,0.7,4.95,12.0,1.5,fill=RGBColor(0xEC,0xF6,0xF4),round_=True)
text(s,1.0,5.13,11.4,1.2,[[("结论： ",15,MINT,True),("内容多样化把对全新部位(人脚)的泛化从“崩溃(r=−0.15)”救成“正常(r≈0.85~0.97)”——",15,INK,False)],
    [("这是全项目最有力的泛化证据；而且脚根本没进训练。",15,INK,False)]],space=1.12)

# ============ 13 群体对照 3x3x1 ============
s=slide(); motif_title(s,"结果④：噪声 OOD 群体对照（3×3×1，36 场景）",13)
why(s,0.7,1.5,12.0,"单张方差大；用 36 个没见过的 3×3×1 场景求均值，才有统计意义的 N2N vs NTN 结论。")
tb=[("模型 / 方法","PSNR","MSSIM","r"),
    ("旧(5x5)  N2N","24.77","0.807","0.908"),("旧(5x5)  NTN","23.88","0.809","0.905"),
    ("新(多被试) N2N","26.81","0.816","0.917"),("新(多被试) NTN","26.91","0.818","0.912")]
yy=2.45
for i,row in enumerate(tb):
    bg=SOFT if i==0 else (RGBColor(0xF7,0xFB,0xFA) if i>=3 else WHITE); box(s,0.7,yy,7.2,0.5,fill=bg,line=LINE)
    for cx,txt in zip([0.9,4.0,5.2,6.5],row):
        text(s,cx,yy+0.09,1.5,0.4,[[(txt,12.5,(NAVY if i==0 else INK),i==0)]])
    yy+=0.5
bullets(s,8.2,2.5,4.6,4.0,[
    [("多被试让 N2N、NTN 都涨（NTN +3dB，更多）。",13.5,INK,False)],
    [("但 NTN 与 N2N 基本打平（差 ±0.1dB）。",13.5,RED,True)],
    [("→ N2N 越强，NTN 越没发挥空间。",13.5,INK,False)],
],sz=13.5,gap=9)

# ============ 14 诚实结论 ============
s=slide(); motif_title(s,"诚实结论：为什么 NTN 没能超过 N2N",14)
bullets(s,0.7,1.6,12.0,3.3,[
    [("结构性原因，不是调参：",15,NAVY,True)],
    [("① 无真值 GT → NTN 用 N2N 当伪 GT(Ĉ=N2N) → 本质“蒸馏 N2N”，天花板就是 N2N；",14.5,INK,False)],
    [("② 你的 N2N(Noise2Noise) 本身是鲁棒强基线，不像论文打败的“过拟合监督去噪器”；",14.5,INK,False)],
    [("③ 试图用小数据削弱 N2N，同时也削弱了 NTN，固定 σ 还帮倒忙（详见附录探索）。",14.5,INK,False)],
],sz=14.5,gap=9)
box(s,0.7,4.85,12.0,1.7,fill=SOFT,line=LINE,round_=True)
text(s,1.0,5.05,11.4,1.4,[[("可写入汇报的真实结论： ",14.5,MINT,True),("①复现了翻译机制(噪声白化,硬证据)；",14.5,INK,False)],
    [("②NTN 泛化与 N2N 相当、不崩不脑补，且 T 与去噪器解耦可即插即用；",14.5,INK,False)],
    [("③真正的泛化提升来自内容多样化；④诚实指出无 GT 时 NTN ≈ N2N 及其原因。",14.5,INK,False)]],space=1.12)

# ============ 15 探索① BSN ============
s=slide(); motif_title(s,"探索①：T 白化后，单帧自监督 BSN 是否可行",15)
why(s,0.7,1.45,12.0,"T 把噪声白化(lag1 0.82→0.13)，正好满足 BSN“噪声空间独立”的前提——验证一条无需 D 的自监督路。")
pic(s,os.path.join(R,'images','bsn_3x3x1','5x5x4_0_npy_0_bsn.png'),0.6,2.3,12.1,2.7)
cap(s,0.6,5.05,12.1,"noisy / BSN(raw) / T(I) / BSN(translated) / NTN=D'(T) / reference （ZS-N2N 单图自监督）")
bullets(s,0.7,5.5,12.0,1.4,[
    [("BSN(raw)≈noisy（去不动相关散斑）；BSN(translated) 明显变好 → T 的白化让单帧自监督去噪可行。",13.5,INK,False)],
    [("但需 T 在目标数据上稳定白化：噪声极强的 3×3×1 上 T 会失稳（见下页）。",13.5,RED,True)],
],sz=13.5,gap=6)

# ============ 16 探索② OOD噪声再测 ============
s=slide(); motif_title(s,"探索②：再测 OOD 噪声，揭示“严重噪声漂移”",16)
why(s,0.7,1.5,12.0,"固定-σ 实验崩了；要把 σ 定准、也要搞清 T 为何在某些 OOD 失稳——所以先测 OOD 噪声。")
tb=[("数据","σ_mad(log1p)"),("训练 5x5 level3","0.14"),("3×3×1 (OOD)","0.64"),("mix 191–226 (OOD)","0.68")]
yy=2.5
for i,row in enumerate(tb):
    bg=SOFT if i==0 else WHITE; box(s,0.7,yy,5.4,0.5,fill=bg,line=LINE)
    for cx,txt in zip([0.95,3.3],row): text(s,cx,yy+0.09,2.2,0.4,[[(txt,13,(NAVY if i==0 else (RED if i>=2 else INK)),i==0 or i>=2)]])
    yy+=0.5
bullets(s,6.5,2.45,6.2,4.0,[
    [("OOD 噪声 ≈0.64–0.68，是训练(0.14)的 4–5 倍，",14,INK,False),("且超出 D' 盲区间上界 0.6。",14,RED,True)],
    [("解释了：固定 σ=0.15 崩（T 要硬减 0.68→0.15）、T 在 3×3×1 失稳（噪声超出训练范围）。",13.5,INK,False)],
    [("启示：若要服务这类超噪 OOD，D' 区间应扩到 ~0.7；单一固定 σ 不适合宽噪声跨度。",13.5,MINT,False)],
],sz=14,gap=9)

# ============ 17 总结 ============
s=slide(NAVY); box(s,0,0,0.18,SH,fill=MINT)
text(s,1.0,0.6,11.4,0.8,[[("总结 & 下一步",30,WHITE,True,TITLEF)]])
bullets(s,1.0,1.8,11.4,3.8,[
    [("✓ 翻译机制复现成功：",16,WHITE,True),("真实散斑被 decorrelate 成近白高斯(硬证据)。",16,ICE,False)],
    [("✓ 内容多样化救活 OOD：",16,WHITE,True),("人脚从崩溃(r=−0.15)到正常(r≈0.9)。",16,ICE,False)],
    [("✓ 新颖应用：",16,WHITE,True),("T 白化 → 单帧自监督 BSN 可行(无需 D/外部 GT)。",16,ICE,False)],
    [("⚠ 诚实边界：",16,MINT,True),("无真值 GT 时 NTN≈N2N(被 N2N 锚定 + N2N 是强基线)。",16,ICE,False)],
    [("下一步：",16,MINT,True),("D' 区间扩到~0.7 服务超噪 OOD；多帧均值目标解耦；稳定 T 的白化。",16,ICE,False)],
],sz=16,gap=13,col=WHITE,marker=MINT)
text(s,1.0,6.7,11.4,0.4,[[("代码 / 实验全记录：github.com/sinyinler/NTN · experiment_log.md",12,MUTE,False)]])

prs.save(OUT); print("saved",OUT,"slides",len(prs.slides._sldIdLst))
