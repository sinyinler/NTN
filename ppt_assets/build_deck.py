# -*- coding: utf-8 -*-
"""构建 NTN 复现汇报 PPT（面向无背景听众）。"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from PIL import Image

A = r'D:\Desktop\NTN\ppt_assets'
OUT = r'D:\Desktop\NTN\NTN_复现汇报.pptx'

NAVY=RGBColor(0x21,0x29,0x5C); BLUE=RGBColor(0x06,0x5A,0x82); TEAL=RGBColor(0x1C,0x72,0x93)
MINT=RGBColor(0x00,0xA8,0x96); INK=RGBColor(0x2A,0x2E,0x35); MUTE=RGBColor(0x6B,0x76,0x82)
WHITE=RGBColor(0xFF,0xFF,0xFF); SOFT=RGBColor(0xF2,0xF5,0xF7); LINE=RGBColor(0xD9,0xE0,0xE5)
RED=RGBColor(0xC0,0x56,0x3A)
TITLEF="Microsoft YaHei"; BODYF="Microsoft YaHei"

prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
BLANK=prs.slide_layouts[6]
SW,SH=13.333,7.5

def slide(bg=WHITE):
    s=prs.slides.add_slide(BLANK)
    r=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,prs.slide_width,prs.slide_height)
    r.fill.solid(); r.fill.fore_color.rgb=bg; r.line.fill.background()
    r.shadow.inherit=False
    s.shapes._spTree.remove(r._element); s.shapes._spTree.insert(2,r._element)
    return s

def box(s,l,t,w,h,fill=None,line=None,round_=False,line_w=1.0):
    shp=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if round_ else MSO_SHAPE.RECTANGLE,
                           Inches(l),Inches(t),Inches(w),Inches(h))
    if fill is None: shp.fill.background()
    else: shp.fill.solid(); shp.fill.fore_color.rgb=fill
    if line is None: shp.line.fill.background()
    else: shp.line.color.rgb=line; shp.line.width=Pt(line_w)
    shp.shadow.inherit=False
    return shp

def text(s,l,t,w,h,runs,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,space=1.05):
    tb=s.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h)); tf=tb.text_frame
    tf.word_wrap=True; tf.vertical_anchor=anchor
    tf.margin_left=0; tf.margin_right=0; tf.margin_top=0; tf.margin_bottom=0
    if isinstance(runs[0],tuple): runs=[runs]
    for i,para in enumerate(runs):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.alignment=align; p.line_spacing=space; p.space_after=Pt(4)
        for (txt,sz,col,bold,*rest) in para:
            r=p.add_run(); r.text=txt; f=r.font
            f.size=Pt(sz); f.color.rgb=col; f.bold=bold; f.name=(rest[0] if rest else BODYF)
    return tb

def motif_title(s,title,num,dark=False):
    box(s,0.6,0.62,0.26,0.26,fill=MINT)
    text(s,1.0,0.5,11.6,0.7,[[(title,30,WHITE if dark else NAVY,True,TITLEF)]])
    text(s,12.2,0.55,0.7,0.5,[[("%02d"%num,14,MUTE,True,TITLEF)]],align=PP_ALIGN.RIGHT)

def pic_fit(s,path,l,t,w,h):
    iw,ih=Image.open(path).size; a=iw/ih; ba=w/h
    if a>ba: nw=w; nh=w/a
    else: nh=h; nw=h*a
    s.shapes.add_picture(path,Inches(l+(w-nw)/2),Inches(t+(h-nh)/2),Inches(nw),Inches(nh))

def bullets(s,l,t,w,h,items,sz=15,gap=8,col=INK,marker=MINT):
    tb=s.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h)); tf=tb.text_frame
    tf.word_wrap=True; tf.margin_left=0; tf.margin_top=0; tf.margin_right=0
    for i,it in enumerate(items):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.line_spacing=1.12; p.space_after=Pt(gap)
        r=p.add_run(); r.text="▪  "; r.font.size=Pt(sz); r.font.color.rgb=marker; r.font.bold=True; r.font.name=BODYF
        if isinstance(it,str): it=[(it,sz,col,False)]
        for (txt,s2,c2,b2) in it:
            r=p.add_run(); r.text=txt; r.font.size=Pt(s2); r.font.color.rgb=c2; r.font.bold=b2; r.font.name=BODYF
    return tb

# ---------- Slide 1: Title ----------
s=slide(NAVY)
box(s,0.0,0.0,0.18,SH,fill=MINT)
text(s,1.1,2.25,11.2,2.0,[
    [("把噪声 “翻译” 成高斯",40,WHITE,True,TITLEF)],
    [("用 NTN 提升 BFI 去噪的泛化性",40,MINT,True,TITLEF)],
],space=1.1)
text(s,1.12,4.5,11.0,1.2,[
    [("复现并改造《Learning to Translate Noise for Robust Image Denoising》",16,RGBColor(0xCA,0xDC,0xFC),False)],
    [("arXiv:2412.04727 → 无真值 GT 的 BFI / Noise2Noise 场景",16,RGBColor(0xCA,0xDC,0xFC),False)],
])
text(s,1.12,6.5,11.0,0.5,[[("宋亚栋  ·  组会汇报  ·  2026-06-18",13,MUTE,False)]])

# ---------- Slide 2: Motivation ----------
s=slide(); motif_title(s,"背景：去噪模型“换个噪声就翻车”",1)
bullets(s,0.7,1.7,7.0,4.8,[
    [("BFI（血流成像）原始图噪声很大，",15,INK,False),("去噪是后续分析的必要前处理。",15,INK,False)],
    [("深度去噪器在“训练时见过的噪声”上很好，",15,INK,False),("一旦换采集设置/噪声分布就明显变差",15,RED,True),("（OOD 泛化差）。",15,INK,False)],
    [("我们已有的 Noise2Noise(N2N) 也一样：",15,INK,False),("只在见过的噪声上强。",15,INK,False)],
    [("目标：",15,INK,True),("让模型对“没见过的噪声”也稳——正是本文要解决的问题。",15,INK,False)],
],sz=15,gap=14)
# 右侧概念图
cx=8.2
box(s,cx,2.0,4.4,1.35,fill=SOFT,line=LINE,round_=True)
text(s,cx+0.25,2.18,3.9,1.0,[[("训练时见过的噪声",13,MUTE,False)],[("去噪效果好  ✔",18,TEAL,True)]],space=1.1)
box(s,cx,3.95,4.4,1.35,fill=SOFT,line=LINE,round_=True)
text(s,cx+0.25,4.13,3.9,1.0,[[("没见过的新噪声 (OOD)",13,MUTE,False)],[("残留 / 过平滑  ✘",18,RED,True)]],space=1.1)
ar=s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW,Inches(cx+1.95),Inches(3.4),Inches(0.5),Inches(0.5))
ar.fill.solid(); ar.fill.fore_color.rgb=MUTE; ar.line.fill.background(); ar.shadow.inherit=False

# ---------- Slide 3: Positioning ----------
s=slide(); motif_title(s,"定位：它和已有方法的区别",2)
cards=[("合成高斯 / 监督去噪","只认高斯噪声；遇到真实噪声直接崩。",TEAL),
       ("N2N / 真实数据训练","受限于训练噪声分布，OOD 仍会过拟合。",BLUE),
       ("NTN（本文）","不直接去未知噪声，先“翻译”成高斯，再用通用高斯去噪器 → 任意噪声都能处理。",MINT)]
cw=3.95; gap=0.28; x0=0.7
for i,(h,b,c) in enumerate(cards):
    x=x0+i*(cw+gap)
    box(s,x,1.9,cw,4.1,fill=SOFT,line=LINE,round_=True)
    box(s,x,1.9,cw,0.12,fill=c)
    text(s,x+0.3,2.25,cw-0.6,1.0,[[(h,18,c,True)]])
    text(s,x+0.3,3.2,cw-0.6,2.6,[[(b,15,INK,False)]],space=1.15)
text(s,0.7,6.35,12,0.5,[[("一句话：把“去各种噪声”的难题，转化为“把各种噪声翻译成同一种高斯”的简单问题。",13.5,MUTE,True)]])

# ---------- Slide 4: Core idea ----------
s=slide(); motif_title(s,"核心思路：与其硬去噪，不如先“翻译”噪声",3)
text(s,0.7,1.6,11.9,1.4,[
    [("关键观察（论文 Fig.1）：",15,NAVY,True),("给真实噪声图再加一点高斯噪声，反而更好去——因为去噪器最擅长高斯。",15,INK,False)],
    [("于是：先用一个小网络 ",15,INK,False),("T 把未知噪声翻译成“空间不相关的高斯噪声”",15,MINT,True),("，再交给专门去高斯的网络 D。",15,INK,False)],
],space=1.15)
pic_fit(s,os.path.join(A,'fig1_teaser.png'),0.7,3.15,11.9,3.4)
text(s,0.7,6.62,11.9,0.4,[[("论文 Fig.1：OOD 噪声图 → 噪声翻译网络 → 高斯噪声图 → 高斯去噪器 → 干净图。",11,MUTE,False)]],align=PP_ALIGN.CENTER)

# ---------- Slide 5: Method ----------
s=slide(); motif_title(s,"方法：T 翻译 + D 去高斯（两阶段）",4)
pic_fit(s,os.path.join(A,'fig3_arch.png'),0.7,1.55,7.4,3.0)
text(s,0.7,4.62,7.4,0.4,[[("论文 Fig.3：训练框架（T 可训练🔥，D 冻结❄）。",11,MUTE,False)]],align=PP_ALIGN.CENTER)
bullets(s,8.35,1.7,4.4,5.0,[
    [("阶段一：",15,NAVY,True),("训练只会去“已知高斯噪声”的去噪器 D（偏科生）。",15,INK,False)],
    [("阶段二：",15,NAVY,True),("冻结 D，训练翻译器 T，使 D(T(I)) ≈ 干净图。",15,INK,False)],
    [("两个损失：",15,NAVY,True),("implicit（去噪结果像干净图）+ explicit（用 Wasserstein 把翻译后噪声逼成高斯：空间域正态、频域 Rayleigh → 更白）。",15,INK,False)],
    [("泛化来源：",15,MINT,True),("D 只认高斯（不会过拟合某种真实噪声），T 负责把任意噪声搬到高斯工作点。",15,INK,False)],
],sz=15,gap=14)

# ---------- Slide 6: Our adaptation ----------
s=slide(); motif_title(s,"我们的改造：无 GT 的 BFI / N2N 版本",5)
bullets(s,0.7,1.65,12.0,2.9,[
    [("没有真干净 GT → ",15,INK,False),("用 N2N 输出当伪干净 Ĉ",15,MINT,True),("（与输入内容对齐，翻译后噪声才“纯”、不误伤血管）。",15,INK,False)],
    [("D' 改成“盲高斯专家”：σ∈[0.08, 0.6]，",15,INK,False),("覆盖实测真实噪声跨度",15,MINT,True),("（log1p 域 level1≈0.43 ~ level4≈0.10，≈1/N），下界远离 0 堵死“恒等捷径”。",15,INK,False)],
    [("explicit loss 后半段才启用（防早期翻译噪声偏差大、训练崩）。",15,INK,False)],
    [("公平实验：N2N 与 NTN 都只用 level 2/3/4 训练，",15,INK,False),("最噪的 level1 留作 OOD 测试。",15,RED,True)],
],sz=15,gap=12)
# pipeline
py=5.15; bw=2.5; bh=1.0; gx=0.55; x=0.95
labels=[("噪声图 I",SOFT,INK),("翻译器 T",MINT,WHITE),("盲高斯 D'",BLUE,WHITE),("干净图",TEAL,WHITE)]
for i,(lab,fc,tc) in enumerate(labels):
    bx=x+i*(bw+gx)
    box(s,bx,py,bw,bh,fill=fc,line=LINE if fc==SOFT else None,round_=True)
    text(s,bx,py,bw,bh,[[(lab,16,tc,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    if i<3:
        ar=s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,Inches(bx+bw+0.05),Inches(py+0.33),Inches(0.45),Inches(0.34))
        ar.fill.solid(); ar.fill.fore_color.rgb=MUTE; ar.line.fill.background(); ar.shadow.inherit=False
text(s,0.95,6.35,11.5,0.4,[[("Ĉ = N2N(I)  作为伪干净目标，同时用作公平对照基线。",12,MUTE,True)]])

# ---------- Slide 7: Quant result ----------
s=slide(); motif_title(s,"结果①：level1（OOD）定量，NTN > N2N",6)
pic_fit(s,os.path.join(A,'chart_psnr.png'),0.6,1.75,6.3,4.6)
bx=7.3
text(s,bx,2.0,5.2,1.0,[[("OOD 泛化增益",16,MUTE,True)]])
text(s,bx,2.5,5.2,1.2,[[("+0.63 dB",54,MINT,True,TITLEF)]])
text(s,bx,3.75,5.2,0.6,[[("NTN 25.64  vs  N2N 25.01 dB",16,INK,True)]])
text(s,bx,4.25,5.2,0.6,[[("SSIM 基本持平（0.758 vs 0.764）",15,INK,False)]])
box(s,bx,5.0,5.2,1.5,fill=SOFT,line=LINE,round_=True)
text(s,bx+0.25,5.18,4.7,1.2,[
    [("评测口径",12,MUTE,True)],
    [("伪GT = 同场景 level4 多帧均值；N2N 与 NTN 同用 level2/3/4 训练、同一 GT 评测，39 场景。",12.5,INK,False)],
],space=1.1)

# ---------- Slide 8: Mechanism ----------
s=slide(); motif_title(s,"结果②：翻译机制成立（硬证据）",7)
pic_fit(s,os.path.join(A,'chart_decorr.png'),0.5,1.65,6.2,3.7)
pic_fit(s,os.path.join(A,'chart_essential.png'),6.85,1.65,6.0,3.7)
bullets(s,0.7,5.55,12.0,1.7,[
    [("真实散斑噪声强空间相关（自相关 0.82）→ 翻译后近白（0.13），谱平坦度 ×9 → ",14.5,INK,False),("确实被翻译成高斯。",14.5,MINT,True)],
    [("把真实噪声直接喂盲高斯 D' 只有 14.2 dB（比噪声还差）；翻译后 20.3 dB（+6 dB）→ ",14.5,INK,False),("翻译是“救命”的一步，且 NTN(20.3) > N2N(19.3)。",14.5,MINT,True)],
],sz=14.5,gap=8)

# ---------- Slide 9: Qualitative far-OOD ----------
s=slide(); motif_title(s,"结果③：远 OOD（3×3 / mix）定性，细节更好",8)
box(s,0.7,1.7,12.0,3.4,fill=SOFT,line=LINE,round_=True)
text(s,0.7,3.0,12.0,0.8,[
    [("［在此插入三联对比图：results/images/ood_compare/comparison/*.png］",15,MUTE,True)],
    [("Input | N2N baseline | NTN(ours)，含血管局部放大",12,MUTE,False)],
],align=PP_ALIGN.CENTER,space=1.2)
bullets(s,0.7,5.35,12.0,1.8,[
    [("在完全不同的采集几何(3×3)、不同场景/格式(mix, 含 .lbf)上：",15,INK,True)],
    [("N2N 明显过平滑、糊；NTN 保留更多沿血管细结构、边缘更锐。",15,INK,False)],
    [("规律：越“远”的 OOD，NTN 相对 N2N 的优势越大 —— 与论文论点一致。",15,MINT,True)],
],sz=15,gap=9)

# ---------- Slide 10: Conclusion ----------
s=slide(NAVY)
box(s,0.0,0.0,0.18,SH,fill=MINT)
text(s,1.0,0.7,11.6,0.8,[[("结论 & 下一步",32,WHITE,True,TITLEF)]])
bullets(s,1.0,2.0,11.4,3.6,[
    [("NTN 的“翻译→去高斯”机制在 BFI 上复现成功：",17,WHITE,True),("真实散斑被 decorrelate 成近白高斯。",17,RGBColor(0xCA,0xDC,0xFC),False)],
    [("OOD 上 NTN 优于同数据训练的 N2N：",17,WHITE,True),("level1 +0.63 dB；远 OOD 视觉更优且更保血管。",17,RGBColor(0xCA,0xDC,0xFC),False)],
    [("下一步（可选）：",17,MINT,True),("把伪干净目标从 N2N 换成多帧均值，进一步保细血管、扩大增益。",17,RGBColor(0xCA,0xDC,0xFC),False)],
],sz=17,gap=18,col=WHITE,marker=MINT)
text(s,1.0,6.6,11.4,0.5,[[("代码 / 实验全记录：github.com/sinyinler/NTN · experiment_log.md",12.5,MUTE,False)]])

prs.save(OUT)
print("saved",OUT,"slides=",len(prs.slides._sldIdLst))
