# -*- coding: utf-8 -*-
"""组会汇报 PPT（重排版）：按"为什么→论文怎么用→流程图→实现细节→失败→修复初成→加数据→调数据量→结论→BSN"。
封面与正文统一浅色风格。所有数值取自 D:\\Desktop\\results 下各 metrics.json，与配图同源。"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from PIL import Image

A = r'D:\Desktop\NTN\ppt_assets'
R = r'D:\Desktop\results'
OUT = r'D:\Desktop\NTN\NTN_组会汇报.pptx'

NAVY=RGBColor(0x21,0x29,0x5C); BLUE=RGBColor(0x06,0x5A,0x82); TEAL=RGBColor(0x1C,0x72,0x93)
MINT=RGBColor(0x00,0xA8,0x96); INK=RGBColor(0x2A,0x2E,0x35); MUTE=RGBColor(0x6B,0x76,0x82)
WHITE=RGBColor(0xFF,0xFF,0xFF); SOFT=RGBColor(0xF2,0xF5,0xF7); LINE=RGBColor(0xD9,0xE0,0xE5)
RED=RGBColor(0xC0,0x56,0x3A); ICE=RGBColor(0xCA,0xDC,0xFC); GOLD=RGBColor(0xC8,0x8A,0x2C)
TITLEF="Microsoft YaHei"; BODYF="Microsoft YaHei"

prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
BLANK=prs.slide_layouts[6]; SW,SH=13.333,7.5

def slide(bg=WHITE):
    s=prs.slides.add_slide(BLANK)
    r=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,prs.slide_width,prs.slide_height)
    r.fill.solid(); r.fill.fore_color.rgb=bg; r.line.fill.background(); r.shadow.inherit=False
    s.shapes._spTree.remove(r._element); s.shapes._spTree.insert(2,r._element)
    return s

def box(s,l,t,w,h,fill=None,line=None,round_=False,lw=1.0):
    shp=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if round_ else MSO_SHAPE.RECTANGLE,Inches(l),Inches(t),Inches(w),Inches(h))
    if fill is None: shp.fill.background()
    else: shp.fill.solid(); shp.fill.fore_color.rgb=fill
    if line is None: shp.line.fill.background()
    else: shp.line.color.rgb=line; shp.line.width=Pt(lw)
    shp.shadow.inherit=False; return shp

def text(s,l,t,w,h,runs,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,space=1.06):
    tb=s.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h)); tf=tb.text_frame
    tf.word_wrap=True; tf.vertical_anchor=anchor
    tf.margin_left=0; tf.margin_right=0; tf.margin_top=0; tf.margin_bottom=0
    if isinstance(runs[0],tuple): runs=[runs]
    for i,para in enumerate(runs):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.alignment=align; p.line_spacing=space; p.space_after=Pt(4)
        for (txt,sz,col,bold,*rest) in para:
            r=p.add_run(); r.text=txt; f=r.font
            f.size=Pt(sz); f.color.rgb=col; f.bold=bold; f.name=(rest[0] if rest else BODYF)
    return tb

def motif_title(s,title,num):
    box(s,0.6,0.5,0.24,0.24,fill=MINT)
    text(s,0.98,0.4,11.0,0.7,[[(title,26,NAVY,True,TITLEF)]])
    text(s,12.2,0.46,0.7,0.5,[[("%02d"%num,13,MUTE,True,TITLEF)]],align=PP_ALIGN.RIGHT)

def pic(s,path,l,t,w,h):
    iw,ih=Image.open(path).size; a=iw/ih; ba=w/h
    nw,nh=(w,w/a) if a>ba else (h*a,h)
    s.shapes.add_picture(path,Inches(l+(w-nw)/2),Inches(t+(h-nh)/2),Inches(nw),Inches(nh))

def bullets(s,l,t,w,h,items,sz=14.5,gap=7,col=INK,marker=MINT):
    tb=s.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h)); tf=tb.text_frame
    tf.word_wrap=True; tf.margin_left=0; tf.margin_top=0; tf.margin_right=0
    for i,it in enumerate(items):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.line_spacing=1.12; p.space_after=Pt(gap)
        r=p.add_run(); r.text="▪  "; r.font.size=Pt(sz); r.font.color.rgb=marker; r.font.bold=True; r.font.name=BODYF
        if isinstance(it,str): it=[(it,sz,col,False)]
        for (txt,s2,c2,b2) in it:
            r=p.add_run(); r.text=txt; r.font.size=Pt(s2); r.font.color.rgb=c2; r.font.bold=b2; r.font.name=BODYF
    return tb

def why(s,l,t,w,txt):
    box(s,l,t,w,0.62,fill=RGBColor(0xEC,0xF6,0xF4),round_=True)
    text(s,l+0.2,t+0.09,w-0.4,0.5,[[("为什么这么试： ",12.5,MINT,True),(txt,12.5,INK,False)]],anchor=MSO_ANCHOR.MIDDLE,space=1.04)

def cap(s,l,t,w,txt):
    text(s,l,t,w,0.35,[[(txt,10.5,MUTE,False)]],align=PP_ALIGN.CENTER)

def table(s,l,t,colw,rows,fs=11,rh=0.46,hl=()):
    tw=sum(colw)
    for i,row in enumerate(rows):
        y=t+i*rh
        bg=SOFT if i==0 else (RGBColor(0xF7,0xFB,0xFA) if i in hl else WHITE)
        box(s,l,y,tw,rh,fill=bg,line=LINE)
        cx=l
        for j,cell in enumerate(row):
            col=NAVY if i==0 else (RED if (cell.strip().startswith("−") or "崩" in cell or "反相关" in cell) else INK)
            text(s,cx+0.12,y+0.09,colw[j]-0.16,0.34,[[(cell,fs,col,i==0)]],anchor=MSO_ANCHOR.MIDDLE); cx+=colw[j]

def fnode(s,l,t,title,sub,c,w=2.15,h=1.0):
    box(s,l,t,w,h,fill=WHITE,line=c,round_=True,lw=1.4)
    text(s,l+0.1,t+0.12,w-0.2,h-0.2,[[(title,11.5,c,True)],[(sub,9.3,INK,False)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE,space=1.04)

def rar(s,l,t):
    a=s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,Inches(l),Inches(t),Inches(0.28),Inches(0.2))
    a.fill.solid(); a.fill.fore_color.rgb=MUTE; a.line.fill.background(); a.shadow.inherit=False

# ==================== 封面（浅色） ====================
s=slide(); box(s,0.6,0.95,0.32,1.7,fill=MINT)
text(s,1.15,1.0,11.4,2.0,[[("把噪声“翻译”成高斯：",34,NAVY,True,TITLEF)],
                          [("NTN 用于 BFI 去噪的复现与泛化性研究",34,NAVY,True,TITLEF)]],space=1.1)
text(s,1.18,3.5,11.0,1.0,[[("复现《Learning to Translate Noise for Robust Image Denoising》(arXiv:2412.04727)，",15,INK,False)],
                          [("并改造到无真值 GT 的 BFI / Noise2Noise 场景",15,INK,False)]],space=1.18)
for i,(t1,t2) in enumerate([("方法定位","噪声 OOD 泛化"),("机制","噪声白化→高斯"),("关键结论","小数据下 NTN 胜 N2N")]):
    x=1.18+i*3.85; box(s,x,4.7,3.55,1.05,fill=SOFT,line=LINE,round_=True)
    text(s,x+0.25,4.82,3.1,0.85,[[(t1,12,MINT,True)],[(t2,13.5,NAVY,True)]],space=1.05)
text(s,1.18,6.55,11,0.5,[[("宋亚栋  ·  组会汇报  ·  2026-06",12.5,MUTE,False)]])

# ==================== 1 为什么用这个方法 / 解决什么 ====================
s=slide(); motif_title(s,"为什么用这个方法：BFI 去噪要解决的问题",1)
bullets(s,0.7,1.55,7.0,4.6,[
    [("BFI（血流成像）原始图噪声大，",15,INK,False),("去噪是后续分析的必要前处理。",15,INK,False)],
    [("深度去噪器在“训练见过的噪声”上很好，",15,INK,False),("换采集/部位/噪声分布就明显变差。",15,RED,True)],
    [("我们的现实约束：",15,NAVY,True),("没有真值 GT（只能用 Noise2Noise 自配对）。",15,INK,False)],
    [("两类 OOD 必须区分：",15,NAVY,True)],
    [("  · 噪声 OOD —— 不同采集/强度（如 3×3 vs 5×5 窗口）",14,INK,False)],
    [("  · 内容 OOD —— 没见过的结构（如鼠脑训练 → 人脚测试）",14,INK,False)],
    [("目标：有限、无 GT 条件下，对 OOD 也稳。",15,MINT,True)],
],sz=15,gap=10)
box(s,8.1,1.9,4.5,1.25,fill=SOFT,line=LINE,round_=True)
text(s,8.35,2.05,4.0,1.0,[[("训练见过的噪声",12.5,MUTE,False)],[("去噪好  ✔",18,TEAL,True)]],space=1.08)
ar=s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW,Inches(10.1),Inches(3.25),Inches(0.45),Inches(0.55))
ar.fill.solid(); ar.fill.fore_color.rgb=MUTE; ar.line.fill.background(); ar.shadow.inherit=False
box(s,8.1,3.95,4.5,1.25,fill=SOFT,line=LINE,round_=True)
text(s,8.35,4.1,4.0,1.0,[[("没见过的噪声 / 结构 (OOD)",12.5,MUTE,False)],[("残留 / 过平滑 / 脑补  ✘",16,RED,True)]],space=1.08)
why(s,8.1,5.55,4.5,"NTN 主打的正是“噪声 OOD”——不直接去未知噪声，而是先翻译成高斯。")

# ==================== 2 论文里怎么用这个方法 ====================
s=slide(); motif_title(s,"论文里怎么用：先翻译，再用通用高斯去噪器",2)
bullets(s,0.7,1.5,5.9,4.4,[
    [("论文场景：真实相机 sRGB 去噪",15,NAVY,True),("（SIDD 等跨相机/ISO）。",14,INK,False)],
    [("痛点：监督去噪器只认训练噪声，",14.5,INK,False),("换噪声分布就崩。",14.5,RED,True)],
    [("做法两段式：",15,NAVY,True)],
    [("① 翻译器 T：把任意未知噪声 → 标准高斯白噪声；",14,INK,False)],
    [("② 高斯去噪器 D：只需会去高斯，固定不动（❄冻结）。",14,INK,False)],
    [("只训 T（🔥），D 即插即用 → 换噪声不必重训去噪器。",14.5,MINT,True)],
    [("效果：在多个 OOD 基准上超过直接训练的去噪器。",14.5,INK,False)],
],sz=14.5,gap=10)
pic(s,os.path.join(A,'fig3_arch.png'),6.8,1.65,6.0,3.5)
cap(s,6.8,5.2,6.0,"论文 Fig.3：翻译器 T 可训练🔥，高斯去噪器 D 冻结❄、即插即用")
why(s,6.8,5.75,6.0,"把“去未知噪声”难题，转成“把噪声标准化 + 去高斯”两个更可控的子问题。")

# ==================== 3 项目流程图（论文风格） ====================
s=slide(); motif_title(s,"项目全流程图",3)
text(s,0.62,1.35,3.0,0.3,[[("训练阶段",12,TEAL,True)]])
top=[("原始散斑 → BFI","log1p 强度域",TEAL),
     ("估计噪声 σ̂","邻帧差 / MAD",BLUE),
     ("训练 N2N","→ 伪干净 Ĉ",NAVY),
     ("训练 D′ 高斯去噪","盲 σ 区间",TEAL),
     ("训练 T 翻译器","隐式+显式损失",MINT)]
xs=[0.55,3.05,5.55,8.05,10.55]
for (l,(ti,su,c)) in zip(xs,top):
    fnode(s,l,1.7,ti,su,c)
for l in xs[:-1]:
    rar(s,l+2.18,2.1)
# 连接：训练得到 T、D′ → 推理
dn=s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW,Inches(11.35),Inches(2.78),Inches(0.45),Inches(1.5))
dn.fill.solid(); dn.fill.fore_color.rgb=MINT; dn.line.fill.background(); dn.shadow.inherit=False
text(s,7.4,3.2,3.8,0.5,[[("训练得到 T、D′，用于推理 →",12,MINT,True)]],align=PP_ALIGN.RIGHT)
text(s,0.62,4.35,3.0,0.3,[[("推理阶段",12,GOLD,True)]])
bot=[("输入 BFI 帧 I","未知噪声",NAVY),
     ("T 白化","→ 高斯噪声",MINT),
     ("D′ 去噪","去标准高斯",TEAL),
     ("仿射校准","去全局尺度偏移",GOLD),
     ("干净 BFI 输出","用于分析",BLUE)]
for (l,(ti,su,c)) in zip(xs,bot):
    fnode(s,l,4.7,ti,su,c)
for l in xs[:-1]:
    rar(s,l+2.18,5.1)
box(s,0.55,6.05,12.2,0.95,fill=SOFT,line=LINE,round_=True)
text(s,0.8,6.16,11.7,0.75,[[("贯穿全程的两条评测线： ",12.5,NAVY,True),
    ("噪声 OOD（3×3 vs 5×5）看 T 是否泛化；内容 OOD（人脚）看数据多样性；",12.5,INK,False)],
    [("指标 PSNR / MSSIM / r 均在 log1p 域、并对输出做仿射校准后计算（去全局尺度偏移，更公平）。",12.5,INK,False)]],space=1.06)

# ==================== 4 方法具体实现（细节） ====================
s=slide(); motif_title(s,"方法实现细节：噪声估计 + T 的两个损失",4)
bullets(s,0.6,1.5,6.5,2.3,[
    [("① 噪声强度怎么估计？",15,NAVY,True)],
    [("  相邻帧差 (f₍k₊₁₎−f_k)/√2 在 log1p 域，",13.5,INK,False),("噪声独立故差分=√2·σ；",13.5,INK,False)],
    [("  用 std / MAD robust 估 σ̂ → 定 D′ 的盲训练区间 [0.02, 0.6]。",13.5,INK,False)],
    [("  （BFI 各层级实测 σ 不同，故 D′ 用一个够宽的区间一次盖全。）",12.5,MUTE,False)],
],sz=13.5,gap=6)
bullets(s,0.6,3.75,6.5,3.4,[
    [("② T 如何把别的噪声学成高斯？两个损失合力：",15,NAVY,True)],
    [("  隐式损失：D′( T(I) ) ≈ 伪干净 Ĉ —— 让翻译结果",13.5,INK,False),("“能被高斯去噪器还原”；",13.5,MINT,True)],
    [("  显式损失（Wasserstein 白化）：",13.5,INK,False)],
    [("    · 空间：残差排序后对齐高斯分位 → 逐像素高斯；",13,INK,False)],
    [("    · 频域：功率谱推平（Rayleigh）→ 去空间相关、白噪声。",13,INK,False)],
    [("  合力 = 翻译后既“可去”又“白且高斯”。",13.5,MINT,True)],
],sz=13.5,gap=6)
pic(s,os.path.join(A,'fig3_arch.png'),6.95,1.55,5.9,2.7)
pic(s,os.path.join(A,'chart_decorr.png'),6.95,4.45,5.9,2.3)
cap(s,6.95,6.78,5.9,"白化生效：邻像素噪声相关 lag-1 由 0.82 → 0.13")

# ==================== 5 第一版失败（诊断）====================
s=slide(); motif_title(s,"第一版直接照搬：效果差、几乎无泛化",5)
text(s,0.7,1.45,12.0,0.7,[[("按论文默认配置直接跑 → 效果差、几乎没有泛化。对照论文逐条诊断，病根在 NTN 翻译设计（不在数据 / N2N —— N2N 本身很好）：",14,INK,False)]])
diag=[("① D′ 高斯区间封顶 0.15","成了恒等捷径，且盖不住真实噪声 0.10~0.43（个别 0.57）。",TEAL),
      ("② 两个损失锚不一致","隐式锚 I2、显式锚糊均值 Ĉ；T(I)−Ĉ 含血管 → 误伤血管。",BLUE),
      ("③ explicit 从第 0 步全程开","论文应训练过半才开；过早强约束扰乱学习。",GOLD),
      ("④ lr 0.01 偏高 + GIBlock 空转","注入初值=0 形同虚设；训练不稳。",RED)]
xs=[0.7,6.75]; ys=[2.35,4.25]
for i,(ti,su,c) in enumerate(diag):
    x=xs[i%2]; y=ys[i//2]
    box(s,x,y,5.85,1.65,fill=SOFT,line=c,round_=True,lw=1.4)
    text(s,x+0.28,y+0.22,5.35,1.25,[[(ti,14.5,c,True)],[(su,12.5,INK,False)]],space=1.12)
box(s,0.7,6.35,12.0,0.78,fill=RGBColor(0xEC,0xF6,0xF4),line=LINE,round_=True)
text(s,1.0,6.5,11.4,0.5,[[("→ 据此逐条修复： ",13.5,MINT,True),
    ("实测噪声定 σ、统一双锚、explicit 延迟开、降 lr、GIBlock 生效（下页见初步成效）。",13.5,INK,False)]],anchor=MSO_ANCHOR.MIDDLE)

# ==================== 6 修复初见成效：level1 OOD（NTN 略优）====================
s=slide(); motif_title(s,"修复后初见成效：OOD 上 NTN 略优于 N2N",6)
pic(s,os.path.join(R,'eval_ood','compare','scene2_frame0.png'),0.55,1.6,7.5,3.65)
cap(s,0.55,5.3,7.5,"留出最噪 level1 当 OOD：noisy / N2N / NTN / 伪GT + 血管放大 —— NTN 放大区细血管比 N2N 更清晰")
table(s,8.25,1.7,[2.5,1.25,1.25],[
    ["level1 OOD (39场景)","PSNR","SSIM"],
    ["noisy","16.42","0.197"],
    ["N2N (lv234)","25.01","0.764"],
    ["NTN (ours)","25.64","0.758"],
],fs=11,rh=0.46,hl=(3,))
bullets(s,8.25,3.95,4.55,3.0,[
    [("修复后第一次正式评测：",13.5,NAVY,True)],
    [("· NTN 25.64 > N2N 25.01（",13,INK,False),("+0.63dB",13,MINT,True),("）。",13,INK,False)],
    [("机制硬证据：",13.5,NAVY,True)],
    [("· 不翻译 D′(I) 仅 14.22dB；D′(T(I))=20.34 > N2N 19.31;",13,INK,False)],
    [("· 噪声空间相关 0.82 → 0.13（白化生效）。",13,INK,False)],
    [("→ 翻译机制成立，OOD 上 NTN 初见优势。",13.5,MINT,True)],
],sz=13,gap=6)

# ==================== 7 增加数据：内容 OOD 被解决 ====================
s=slide(); motif_title(s,"增加多被试数据：内容 OOD（人脚）被救回",7)
pic(s,os.path.join(R,'metrics_foot5','compare_gray_color.png'),0.55,1.55,12.2,3.1)
cap(s,0.55,4.66,12.2,"多被试（脑/腿/手/脚）训练后，再测没见过的人脚：NTN 结构恢复正常")
table(s,0.6,5.25,[3.0,1.7,1.7],[
    ["人脚 内容OOD（NTN）","旧:单一数据","新:多被试"],
    ["与真值相关 r","−0.15","0.84"],
    ["与真值 PSNR","2.98","24.66"],
],fs=11,rh=0.44)
bullets(s,7.6,5.3,5.2,2.0,[
    [("关键发现：",14.5,NAVY,True)],
    [("· 内容 OOD 不是 NTN 的活，靠",13,INK,False),("数据多样性",13,MINT,True),("解决；",13,INK,False)],
    [("· 36 个 OOD 场景上 NTN−N2N = −0.02dB（持平）;",13,INK,False)],
    [("· 数据充足时 N2N 也很强，NTN 没拉开。",13,INK,False)],
],sz=13,gap=5)

# ==================== 8 再次调整数据量：limited supervision 优势 ====================
s=slide(); motif_title(s,"再调数据量：数据越少，NTN 优势越大",8)
pic(s,os.path.join(R,'eval_ood_scenes_small_blind_affine','vis','197.png'),0.55,1.5,7.4,3.7)
cap(s,0.55,5.22,7.4,"小数据 + 盲 σ，未见过的 3×3 OOD（场景197）：NTN 比 N2N 高 1.5dB")
table(s,8.15,1.7,[2.55,2.15],[
    ["36 OOD·仿射校准","NTN − N2N"],
    ["小数据(level3×5场景)","+0.85 dB"],
    ["多被试(全量)","−0.02 dB"],
],fs=11.5,rh=0.5)
bullets(s,8.15,3.55,4.6,2.6,[
    [("小数据：",13.5,NAVY,True),("N2N 32.59 → NTN 33.45",13.5,INK,False)],
    [("                r 0.864 → 0.887",12.5,MUTE,False)],
    [("多被试：",13.5,NAVY,True),("34.73 ≈ 34.71（持平）",13.5,INK,False)],
    [("→ 正是论文 “limited supervision”",13.5,MINT,True)],
    [("   下的鲁棒泛化，在 BFI 上验证成立。",13.5,MINT,True)],
],sz=13.5,gap=5)
why(s,8.15,6.0,4.6,"数据少时 N2N 易过拟合、泛化弱，NTN 的翻译+鲁棒 D′ 才显出优势。")

# ==================== 9 结论 ====================
s=slide(); motif_title(s,"结论",9)
bullets(s,0.7,1.65,12.0,4.4,[
    [("① 机制复现成功：",16,NAVY,True),("T 把相关噪声白化成高斯（lag-1 0.82→0.13），硬证据。",16,INK,False)],
    [("② 优势取决于数据量：",16,NAVY,True),("小数据 NTN > N2N（+0.85dB、r+0.023）；数据充足时持平。",16,INK,False)],
    [("③ 两类 OOD 各有归属：",16,NAVY,True),("噪声 OOD 靠 NTN 翻译；内容 OOD 靠数据多样性。",16,INK,False)],
    [("④ 评测/部署需一次仿射校准：",16,NAVY,True),("去全局尺度偏移；r、SSIM 不受影响、结构本就对。",16,INK,False)],
    [("⑤ 额外价值：",16,NAVY,True),("T 与去噪器解耦、即插即用；白化后单图 BSN 才变得可行。",16,INK,False)],
],sz=16,gap=24)
box(s,0.7,5.65,12.0,1.15,fill=RGBColor(0xEC,0xF6,0xF4),line=LINE,round_=True)
text(s,1.0,5.84,11.4,0.85,[[("一句话： ",15,MINT,True),
    ("在无 GT 的 BFI 上，NTN 的噪声泛化是真实有效的——尤其在数据有限时胜过 N2N；其价值是“训练一次跨噪声部署”，而非一味追求高 dB。",14.5,INK,False)]],space=1.1)

# ==================== 10 后续探索：单图 BSN ====================
s=slide(); motif_title(s,"后续探索：白化后单张图 BSN 可行",10)
pic(s,os.path.join(R,'images','bsn_3x3x1_real','3x3x1_0_npy_0_bsn.png'),0.55,1.5,12.2,2.9)
cap(s,0.55,4.42,12.2,"noisy / BSN(raw) / T(I)翻译图 / BSN(translated) / NTN=D′(T) / reference —— 原始散斑上 BSN 去不动，翻译白化后 BSN 生效")
table(s,0.6,5.05,[2.6,1.5,1.5,1.5],[
    ["对同一真值","PSNR","MSSIM","r"],
    ["raw 输入","30.76","0.759","0.833"],
    ["T(I) 白化图","14.15","0.719","0.238"],
    ["N2N / NTN","33.26 / 33.88","0.873 / 0.875","0.902 / 0.900"],
],fs=10.5,rh=0.42)
bullets(s,8.4,5.1,4.4,2.1,[
    [("思路：",14.5,NAVY,True),("BSN 靠“邻像素噪声独立”。",13,INK,False)],
    [("原始散斑空间相关 → BSN 失效；",13,RED,True)],
    [("T 白化后(r→0.24)邻像素近独立 → 单图自监督 BSN 可去噪。",13,INK,False)],
    [("意义：无需配对帧，单张图也能去噪。",13,MINT,True)],
],sz=13,gap=5)

prs.save(OUT)
print("saved",OUT,"slides",len(prs.slides._sldIdLst))
